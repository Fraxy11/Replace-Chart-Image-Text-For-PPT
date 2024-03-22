import os
from fastapi import FastAPI
from fastapi.responses import FileResponse
from pptx import Presentation
from pptx.chart.data import CategoryChartData
from python_pptx_text_replacer import TextReplacer
from pptx_replace import replace_picture
from typing import Any, Dict
from datetime import datetime
import calendar
import base64


app = FastAPI()

BASE_DIR = os.path.dirname(os.path.abspath(__file__))
DIR1 = os.path.join(BASE_DIR, "result")
DIR2 = os.path.join(BASE_DIR, "template")
DIR3 = os.path.join(BASE_DIR,  "image")

def replace_chart_with_data(slide, chart_index, chart_data):
        chart_count = 0
        for shape in slide.shapes:
            if shape.has_chart:
                chart_count += 1
                if chart_count == chart_index + 1:  
                    chart = shape.chart
                    chart.replace_data(chart_data)
                    print(f"data:{chart_data}")


                    print(f"Chart with index {chart_index} found and replaced successfully on the specified slide.")
                    return
            print(f"Chart with index {chart_index} not found on the specified slide.")
            
@app.post('/generate')
def generate(request: Dict[Any, Any]):
    print(f"{request}")
    
    results = request
    result = results['result']
    each_day = result['each_day_count']
    early_date = result['earliest_date']
    last_date = result['latest_date']

    file = "template.pptx"
    template_path = os.path.join(DIR2, file)

    month = datetime.strptime(early_date, "%Y-%m-%d")
    bulan = month.month
    nama_bulan = calendar.month_name[bulan]


    prs = Presentation(template_path)

    hari = []
    total = []

    for day_count_dict in each_day:
        for date, count in day_count_dict.items():
            hari.append(date)
            total.append(count)
    # print(hari)
    # print(total)
    dates = [datetime.strptime(date, "%Y-%m-%d").date() for date in hari]
    slide_index = 1
    chart_index_to_replace = 0
    new_new_chart_data = CategoryChartData()
    new_new_chart_data.categories = dates
    new_new_chart_data.add_series('',total )

    for i, slide in enumerate(prs.slides):
        if i == slide_index:
            replace_chart_with_data(slide, chart_index_to_replace, new_new_chart_data)
        break 
    for i, slide in enumerate(prs.slides):
        if i == slide_index:
            replace_chart_with_data(slide, chart_index_to_replace, new_new_chart_data)
    
    # REPLACE IMAGE
    image = 'img.png'
    img= result['sna']['image']
    SAVE_F = os.path.join(DIR3, image )

    img_bytes = base64.b64decode(img)
    with open(SAVE_F, "wb") as fh:
        fh.write(img_bytes)
    slide = prs.slides[3]

    replace_picture(slide, SAVE_F, auto_reshape=True, pic_number=0, )
        
    img_data1 = result['list_of_images'][0]

    img_bytes = base64.b64decode(img_data1)
    with open(SAVE_F, "wb") as fh:
        fh.write(img_bytes)
    slide = prs.slides[2]

    replace_picture(slide, SAVE_F, auto_reshape=True, pic_number=2, )

    img_data2 = result['list_of_images'][1]

    img_bytes = base64.b64decode(img_data2)
    with open(SAVE_F, "wb") as fh:
        fh.write(img_bytes)
    slide = prs.slides[2]

    replace_picture(slide, SAVE_F, auto_reshape=True, pic_number=3, )
    
    img_data3 = result['list_of_images'][2]

    img_bytes = base64.b64decode(img_data3)
    with open(SAVE_F, "wb") as fh:
        fh.write(img_bytes)
    slide = prs.slides[2]

    replace_picture(slide, SAVE_F, auto_reshape=True, pic_number=4, )
    
    new_filename = "result.pptx"
    SAVE_FS = os.path.join(DIR1, new_filename)
    prs.save(SAVE_F)

    trend = result['trend_analysis']
    count = result['total_count']
    topic = result['topics']
    cluster = result['sna']['clusters'][0]
    sumary = cluster['summary']
    statis = result['sna']['statistics']
    catatan = result['sna']['summary']
    platform = result['platform_count']

    def format_angka(angka):
        angka_format = "{:,.0f}".format(angka)
        angka_format = angka_format.replace(',', '.')
        return angka_format

    replace = TextReplacer(SAVE_FS, slides='', tables=True, charts=False, textframes=True)
    replace.replace_text([
        ('1 – 15 Januari', early_date[8:10]+"–"+last_date[8:10]+" "+nama_bulan),
        ('1 - 15 Januari', early_date[8:10]+"-"+last_date[8:10]+" "+nama_bulan),
        ('Perhatian terhadap IKN cenderung tidak meningkat namun sempat mengalami lonjakan karena isu-isu tertentu.', trend),
        ('36.343', format_angka(count)),
        ('Pasca debat, Prabowo dikritik karena kepemilikan tanah di IKN.', str(topic[0])if isinstance(topic, list) and len(topic) > 0 else ''),
        ('Pemerintah dikritik karena ingin menggelontorkan triliunan uang untuk pembangunan IKN.',str(topic[1])if isinstance(topic, list) and len(topic) > 1 else ''),
        ('Netizen soroti isu Djarum dan Wings Group hengkang dari konsorsium IKN.', str(topic[2])if isinstance(topic, list) and len(topic) > 2 else ''),
        # STATISTIC
        ('7.961', str(statis['account_count'])),
        ('100',str(statis['hashtag_count'])),
        ('22.542',str(statis['activity_count'])),
        # PRO IKN
        ('Kelompok ini cenderung merupakan akun-akun pro pemerintah dan pro Prabowo.',str(sumary[0])),
        ('Kelompok ini angkat keberhasilan pemerintah mendapat investasi 7 Triliun setelah kunjungan ke negara-negara Asean.',str(sumary[1])),
        ('Kelompok pro Prabowo kritik Anies yang dianggap sikapnya kini mulai tidak konsisten terhadap IKN yang sebelumnya aktif menolak.',str(sumary[2])),
        ('Kelompok ini klarifikasi isu adanya investor yang mundur dari IKN.',str(sumary[3])),
        ('Kelompok ini menunjukkan dampak positif IKN terhadap daerah sekitarnya yang akan ikut maju.',str(sumary[4])),
        # # CATATAN
        ('Kelompok pro IKN aktif angkat keberhasilan Jokowi mendapat investasi 7 triliun setelah mengadakan kunjungan ke ASEAN, terutama dari Brunei.',str(catatan[0])),
        ('Kelompok kontra IKN masih terus menyindir IKN karena menunjukkan sikap pemerintah yang hanya ingin menguntungkan diri sendiri, bukan pro rakyat.',str(catatan[1])),
        ('Kontra IKN juga angkat isu adanya investor IKN yang keluar dari konsorsium.',str(catatan[2])),
        ('26.953 Data (74',str(platform[0]['twitter']['total'])+" Data ("+str(platform[0]['twitter']['percentage']*100)),
        ('586 Data (1.6',str(platform[3]['facebook']['total'])+" Data ("+str(platform[3]['facebook']['percentage']*100)),
        ('7.351 Data (20',str(platform[1]['youtube']['total'])+" Data ("+str(platform[1]['youtube']['percentage']*100)),
        ('726 Data (1.99',str(platform[2]['instagram']['total'])+" Data ("+str(platform[2]['instagram']['percentage']*100)),
        ('729 Data (2',str(platform[4]['tiktok']['total'])+" Data ("+str(platform[4]['tiktok']['percentage']*100)),

    ])

    replace.write_presentation_to_file(SAVE_F) 

    file_path = os.path.join(DIR1, new_filename)
    return FileResponse(path=file_path, filename=new_filename)
